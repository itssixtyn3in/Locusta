import random
import zipfile
from pathlib import Path

from docx import Document
from docx.shared import RGBColor, Pt, Inches
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls

from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.comments import Comment

from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGBColor

from fpdf import FPDF

from PIL import Image
from stegano import lsb


WHITE_RGB = (255, 255, 255)
WHITE_HEX = "FFFFFF"

WHITE_DOCX = RGBColor(*WHITE_RGB)
WHITE_PPTX = PptxRGBColor(*WHITE_RGB)


# ------------------------------------------------
# HELPERS
# ------------------------------------------------

def generate_account_numbers(
    count=5,
    min_digits=8
):

    lower = 10 ** (min_digits - 1)

    upper = (10 ** min_digits) - 1

    return [
        random.randint(lower, upper)
        for _ in range(count)
    ]


def prepare_generation(output_folder):

    output_path = Path(output_folder)

    output_path.mkdir(
        parents=True,
        exist_ok=True
    )

    paid = generate_account_numbers()

    unpaid = generate_account_numbers()

    return paid, unpaid


# ------------------------------------------------
# CONTENT HELPERS
# ------------------------------------------------

def _write_docx_content(
    doc,
    paid,
    unpaid
):

    doc.add_heading(
        "Paid Accounts",
        level=1
    )

    for acc in paid:

        doc.add_paragraph(
            str(acc),
            style="List Bullet"
        )

    doc.add_heading(
        "Unpaid Accounts",
        level=1
    )

    for acc in unpaid:

        doc.add_paragraph(
            str(acc),
            style="List Bullet"
        )


def _write_xlsx_content(
    ws,
    paid,
    unpaid
):

    ws.append(["Type", "Number"])

    for acc in paid:

        ws.append(["Paid", acc])

    ws.append(["---", "---"])

    for acc in unpaid:

        ws.append(["Unpaid", acc])


def _write_pptx_content(
    slide,
    paid,
    unpaid
):

    slide.shapes.title.text = (
        "Account Status"
    )

    paid_str = "\n".join(
        map(str, paid)
    )

    unpaid_str = "\n".join(
        map(str, unpaid)
    )

    tf = slide.placeholders[1].text_frame

    tf.text = (
        f"Paid Accounts:\n{paid_str}\n\n"
        f"Unpaid Accounts:\n{unpaid_str}"
    )


# ------------------------------------------------
# OOXML ORPHAN PART
# ------------------------------------------------

def _inject_ooxml_orphan_part(
    zip_path,
    prompt,
    base_path
):

    with zipfile.ZipFile(
        zip_path,
        "a"
    ) as zf:

        orphan = (
            f'<?xml version="1.0"?>'
            f'<prompt>{prompt}</prompt>'
        ).encode("utf-8")

        zf.writestr(
            f"{base_path}/z_unreferenced_prompt.xml",
            orphan
        )


# ------------------------------------------------
# IMAGE STEGO
# ------------------------------------------------

def _run_image_stego(
    folder,
    filename,
    prompt
):

    temp_dir = Path(folder)

    base_img = (
        temp_dir /
        f"{filename}_base.png"
    )

    secret_img = (
        temp_dir /
        f"{filename}_secret.png"
    )

    Image.new(
        "RGB",
        (60, 20),
        color="white"
    ).save(base_img)

    hidden = lsb.hide(
        str(base_img),
        prompt
    )

    hidden.save(secret_img)

    if base_img.exists():

        base_img.unlink()

    return secret_img


# ------------------------------------------------
# VERIFICATION HELPERS
# ------------------------------------------------

def verify_hidden_text(path, prompt):

    print("\n[VERIFY] Hidden Text")

    print(
        "[+] Prompt inserted as white-colored text"
    )

    print(
        "[+] Select all text inside Office "
        "to reveal hidden content"
    )

    print(f"[+] File: {path}")


def verify_header_footer(path, prompt):

    print("\n[VERIFY] Header/Footer")

    print(
        "[+] Prompt embedded in "
        "document header/footer"
    )

    print(
        "[+] Open Header/Footer editing mode"
    )

    print(f"[+] File: {path}")


def verify_off_page(path, prompt):

    print("\n[VERIFY] Off Page")

    print(
        "[+] Prompt stored in floating "
        "object outside render bounds"
    )

    print(
        "[+] Inspect Selection Pane or XML"
    )

    print(f"[+] File: {path}")


def verify_orphan_stream(path, prompt):

    print("\n[VERIFY] Orphan Stream")

    print(
        "[+] Prompt stored in "
        "unreferenced OOXML object"
    )

    print(
        "[+] Rename file to .zip "
        "and inspect XML"
    )

    if str(path).endswith(".docx"):

        print(
            "[+] Path: word/z_unreferenced_prompt.xml"
        )

    elif str(path).endswith(".xlsx"):

        print(
            "[+] Path: xl/z_unreferenced_prompt.xml"
        )

    print(f"[+] File: {path}")


def verify_comment(path, prompt):

    print("\n[VERIFY] Hidden Comment")

    print(
        "[+] Prompt stored inside comments"
    )

    print(
        "[+] Open Comments Pane in Office"
    )

    print(f"[+] File: {path}")


def verify_style_stego(path, prompt):

    print("\n[VERIFY] Style Steganography")

    print(
        "[+] Prompt encoded in style names"
    )

    print(
        "[+] Inspect styles.xml "
        "inside OOXML archive"
    )

    print(f"[+] File: {path}")


def verify_image_stego(path, prompt):

    print("\n[VERIFY] Image Steganography")

    print(
        "[+] Prompt hidden inside embedded PNG"
    )

    print(
        "[+] Extract image and decode with:"
    )

    print(
        "    stegano.lsb.reveal(image.png)"
    )

    print(f"[+] File: {path}")


# ------------------------------------------------
# HIDDEN TEXT
# ------------------------------------------------

def create_docx_hidden_text(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    p = doc.add_paragraph()

    run = p.add_run(prompt)

    run.font.color.rgb = WHITE_DOCX

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc.save(path)

    return path


def create_xlsx_hidden_text(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    wb = Workbook()

    ws = wb.active

    _write_xlsx_content(
        ws,
        paid,
        unpaid
    )

    cell = ws.cell(
        row=len(paid) + 2,
        column=1,
        value=prompt
    )

    cell.font = Font(
        color=WHITE_HEX
    )

    path = (
        Path(folder) /
        f"{filename}.xlsx"
    )

    wb.save(path)

    return path


def create_pptx_hidden_text(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    prs = Presentation()

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    _write_pptx_content(
        slide,
        paid,
        unpaid
    )

    txbox = slide.shapes.add_textbox(
        PptxInches(0.5),
        PptxInches(6),
        PptxInches(9),
        PptxInches(1)
    )

    p = txbox.text_frame.add_paragraph()

    p.text = prompt

    p.font.color.rgb = WHITE_PPTX

    path = (
        Path(folder) /
        f"{filename}.pptx"
    )

    prs.save(path)

    return path


# ------------------------------------------------
# HEADER FOOTER
# ------------------------------------------------

def create_docx_hf(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    section = doc.sections[0]

    header = section.header

    p = (
        header.paragraphs[0]
        if header.paragraphs
        else header.add_paragraph()
    )

    run = p.add_run(prompt)

    run.font.color.rgb = WHITE_DOCX

    run.font.size = Pt(1)

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc.save(path)

    return path


def create_xlsx_hf(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    wb = Workbook()

    ws = wb.active

    _write_xlsx_content(
        ws,
        paid,
        unpaid
    )

    ws.oddHeader.center.text = (
        f"&1&KFFFFFF{prompt}"
    )

    path = (
        Path(folder) /
        f"{filename}.xlsx"
    )

    wb.save(path)

    return path


# ------------------------------------------------
# OFF PAGE
# ------------------------------------------------

def create_docx_off_page(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    p = doc.add_paragraph()

    xml = rf"""
    <w:pict {nsdecls('w', 'v')}>
      <v:shape
        id="OffPageText"
        style="
          position:absolute;
          margin-left:-5000pt;
          margin-top:-5000pt;
          width:1pt;
          height:1pt;
          z-index:-1;
        "
        fillcolor="white"
        stroked="f">

        <v:textbox inset="0,0,0,0">
          <w:txbxContent>
            <w:p>
              <w:r>
                <w:t>{prompt}</w:t>
              </w:r>
            </w:p>
          </w:txbxContent>
        </v:textbox>

      </v:shape>
    </w:pict>
    """

    p._element.append(
        parse_xml(xml)
    )

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc.save(path)

    return path


def create_pptx_off_page(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    prs = Presentation()

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    _write_pptx_content(
        slide,
        paid,
        unpaid
    )

    txbox = slide.shapes.add_textbox(
        PptxInches(50),
        PptxInches(50),
        PptxInches(1),
        PptxInches(1)
    )

    tf = txbox.text_frame

    tf.text = prompt

    p = tf.paragraphs[0]

    p.font.size = Pt(1)

    p.font.color.rgb = WHITE_PPTX

    path = (
        Path(folder) /
        f"{filename}.pptx"
    )

    prs.save(path)

    return path


# ------------------------------------------------
# ORPHAN STREAM
# ------------------------------------------------

def create_docx_orphan_stream(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    doc.save(path)

    _inject_ooxml_orphan_part(
        path,
        prompt,
        "word"
    )

    return path


def create_xlsx_orphan_stream(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    path = (
        Path(folder) /
        f"{filename}.xlsx"
    )

    wb = Workbook()

    ws = wb.active

    _write_xlsx_content(
        ws,
        paid,
        unpaid
    )

    wb.save(path)

    _inject_ooxml_orphan_part(
        path,
        prompt,
        "xl"
    )

    return path


# ------------------------------------------------
# HIDDEN COMMENT
# ------------------------------------------------

def create_docx_comment(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    comment_p = doc.add_paragraph()

    run = comment_p.add_run(
        f"[COMMENT]: {prompt}"
    )

    run.font.hidden = True

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc.save(path)

    return path


def create_xlsx_comment(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    wb = Workbook()

    ws = wb.active

    _write_xlsx_content(
        ws,
        paid,
        unpaid
    )

    ws["A1"].comment = Comment(
        prompt,
        "System"
    )

    path = (
        Path(folder) /
        f"{filename}.xlsx"
    )

    wb.save(path)

    return path


# ------------------------------------------------
# STYLE STEGO
# ------------------------------------------------

def create_docx_style_stego(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    binary_prompt = "".join(
        format(ord(c), "08b")
        for c in prompt
    )

    for i, bit in enumerate(binary_prompt):

        style = doc.styles.add_style(
            f"Stego_{i}_{bit}",
            WD_STYLE_TYPE.CHARACTER
        )

        style.font.hidden = True

    path = (
        Path(folder) /
        f"{filename}.docx"
    )

    doc.save(path)

    return path


# ------------------------------------------------
# IMAGE STEGO
# ------------------------------------------------

def create_docx_image_stego(
    paid,
    unpaid,
    folder,
    filename,
    prompt
):

    doc = Document()

    _write_docx_content(
        doc,
        paid,
        unpaid
    )

    secret = _run_image_stego(
        folder,
        filename,
        prompt
    )

    try:

        doc.add_picture(
            str(secret),
            width=Inches(0.1)
        )

        path = (
            Path(folder) /
            f"{filename}.docx"
        )

        doc.save(path)

        return path

    finally:

        if secret.exists():

            secret.unlink()


# ------------------------------------------------
# TECHNIQUES
# ------------------------------------------------

TEST_CASES = {

    "hidden_text": {

        "description":
            "Inject hidden white text",

        "generators": {

            "docx":
                create_docx_hidden_text,

            "xlsx":
                create_xlsx_hidden_text,

            "pptx":
                create_pptx_hidden_text
        },

        "verify":
            verify_hidden_text
    },

    "header_footer": {

        "description":
            "Hide prompts inside headers and footers",

        "generators": {

            "docx":
                create_docx_hf,

            "xlsx":
                create_xlsx_hf
        },

        "verify":
            verify_header_footer
    },

    "off_page": {

        "description":
            "Place prompts outside visible page boundaries",

        "generators": {

            "docx":
                create_docx_off_page,

            "pptx":
                create_pptx_off_page
        },

        "verify":
            verify_off_page
    },

    "orphan_stream": {

        "description":
            "Inject orphan OOXML streams and hidden XML",

        "generators": {

            "docx":
                create_docx_orphan_stream,

            "xlsx":
                create_xlsx_orphan_stream
        },

        "verify":
            verify_orphan_stream
    },

    "hidden_comment": {

        "description":
            "Hide prompts inside comments and metadata",

        "generators": {

            "docx":
                create_docx_comment,

            "xlsx":
                create_xlsx_comment
        },

        "verify":
            verify_comment
    },

    "style_steganography": {

        "description":
            "Encode prompts into style metadata",

        "generators": {

            "docx":
                create_docx_style_stego
        },

        "verify":
            verify_style_stego
    },

    "image_steganography": {

        "description":
            "Hide prompts inside embedded images",

        "generators": {

            "docx":
                create_docx_image_stego
        },

        "verify":
            verify_image_stego
    }
}


# ------------------------------------------------
# MAIN GENERATOR
# ------------------------------------------------

def generate_payload(
    technique,
    prompt,
    output_folder="generated",
    filename="payload"
):

    if technique not in TEST_CASES:

        print(
            f"[-] Unknown technique: "
            f"{technique}"
        )

        return

    config = TEST_CASES[technique]

    paid, unpaid = prepare_generation(
        output_folder
    )

    print(
        f"\n[*] Running generation: "
        f"{technique}"
    )

    print(
        f"[+] Technique Description: "
        f"{config['description']}"
    )

    generated = []

    verify_func = config.get(
        "verify"
    )

    verification_printed = False

    for ext, generator in config[
        "generators"
    ].items():

        try:

            path = generator(
                paid,
                unpaid,
                output_folder,
                filename,
                prompt
            )

            generated.append(path)

            print(
                f"[+] Generated: "
                f"{path}"
            )

            # ---------------------------------
            # VERIFY ONLY ONCE
            # ---------------------------------

            if (
                verify_func and
                not verification_printed
            ):

                verify_func(
                    path,
                    prompt
                )

                verification_printed = True

        except Exception as e:

            print(
                f"[-] Failed {ext}: "
                f"{e}"
            )

    print(
        f"\n[+] Total Generated Files: "
        f"{len(generated)}"
    )

    return generated